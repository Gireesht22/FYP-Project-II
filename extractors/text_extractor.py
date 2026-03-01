"""
Text Extractor Module
Supports: PDF, DOCX, TXT, CSV, XLSX, PPTX, HTML, JSON, XML, RTF, MD
No external APIs - all local processing
"""

import os
import csv
import json
import zipfile
import io
import re
import html
from pathlib import Path

# ── PDF ──────────────────────────────────────────────────────────────────────
def extract_pdf(filepath):
    try:
        import pdfplumber
        text_parts = []
        meta = {}
        with pdfplumber.open(filepath) as pdf:
            meta = {
                "pages": len(pdf.pages),
                "metadata": pdf.metadata or {}
            }
            for i, page in enumerate(pdf.pages):
                t = page.extract_text()
                if t:
                    text_parts.append(f"--- Page {i+1} ---\n{t.strip()}")
                # also grab tables
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        row_text = " | ".join(str(c) if c else "" for c in row)
                        text_parts.append(row_text)
        return {
            "text": "\n\n".join(text_parts),
            "metadata": meta,
            "word_count": count_words("\n".join(text_parts))
        }
    except Exception as e:
        return {"text": "", "error": str(e), "metadata": {}}


# ── DOCX ─────────────────────────────────────────────────────────────────────
def extract_docx(filepath):
    try:
        from docx import Document
        doc = Document(filepath)
        sections = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name if para.style else "Normal"
                sections.append({"style": style, "text": para.text.strip()})
        
        # Tables
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_texts.append(" | ".join(cells))
        
        full_text = "\n".join(s["text"] for s in sections)
        if table_texts:
            full_text += "\n\n[Tables]\n" + "\n".join(table_texts)
        
        props = doc.core_properties
        meta = {
            "author": props.author or "",
            "title": props.title or "",
            "created": str(props.created) if props.created else "",
            "paragraphs": len(sections),
            "tables": len(doc.tables)
        }
        return {
            "text": full_text,
            "metadata": meta,
            "word_count": count_words(full_text)
        }
    except Exception as e:
        return {"text": "", "error": str(e), "metadata": {}}


# ── XLSX ─────────────────────────────────────────────────────────────────────
def extract_xlsx(filepath):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, data_only=True)
        all_text = []
        sheets_info = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                clean = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in clean):
                    rows_text.append(" | ".join(clean))
            
            sheet_text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text)
            all_text.append(sheet_text)
            sheets_info.append({
                "name": sheet_name,
                "rows": ws.max_row,
                "cols": ws.max_column
            })
        
        full_text = "\n\n".join(all_text)
        return {
            "text": full_text,
            "metadata": {"sheets": sheets_info},
            "word_count": count_words(full_text)
        }
    except Exception as e:
        return {"text": "", "error": str(e), "metadata": {}}


# ── PPTX ─────────────────────────────────────────────────────────────────────
def extract_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides_text = []
        
        for i, slide in enumerate(prs.slides):
            parts = [f"--- Slide {i+1} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            if len(parts) > 1:
                slides_text.append("\n".join(parts))
        
        full_text = "\n\n".join(slides_text)
        return {
            "text": full_text,
            "metadata": {"slides": len(prs.slides)},
            "word_count": count_words(full_text)
        }
    except Exception as e:
        return {"text": "", "error": str(e), "metadata": {}}


# ── CSV ───────────────────────────────────────────────────────────────────────
def extract_csv(filepath):
    try:
        import chardet
        with open(filepath, 'rb') as f:
            raw = f.read()
        enc = chardet.detect(raw)['encoding'] or 'utf-8'
        
        rows = []
        with open(filepath, encoding=enc, errors='replace') as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" | ".join(row))
        
        full_text = "\n".join(rows)
        return {
            "text": full_text,
            "metadata": {"rows": len(rows), "encoding": enc},
            "word_count": count_words(full_text)
        }
    except Exception as e:
        return {"text": "", "error": str(e), "metadata": {}}


# ── TXT / MD / RTF ────────────────────────────────────────────────────────────
def extract_txt(filepath):
    try:
        import chardet
        with open(filepath, 'rb') as f:
            raw = f.read()
        enc = chardet.detect(raw)['encoding'] or 'utf-8'
        text = raw.decode(enc, errors='replace')
        return {
            "text": text,
            "metadata": {"encoding": enc},
            "word_count": count_words(text)
        }
    except Exception as e:
        return {"text": "", "error": str(e), "metadata": {}}


# ── HTML ──────────────────────────────────────────────────────────────────────
def extract_html(filepath):
    try:
        from bs4 import BeautifulSoup
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        soup = BeautifulSoup(content, 'lxml')
        for tag in soup(["script", "style", "meta", "link"]):
            tag.decompose()
        text = soup.get_text(separator='\n')
        text = re.sub(r'\n{3,}', '\n\n', text).strip()
        title = soup.title.string if soup.title else ""
        return {
            "text": text,
            "metadata": {"title": title},
            "word_count": count_words(text)
        }
    except Exception as e:
        return {"text": "", "error": str(e), "metadata": {}}


# ── JSON ──────────────────────────────────────────────────────────────────────
def extract_json(filepath):
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            data = json.load(f)
        text = json.dumps(data, indent=2)
        flat_text = flatten_json(data)
        return {
            "text": flat_text,
            "raw": text,
            "metadata": {"type": type(data).__name__},
            "word_count": count_words(flat_text)
        }
    except Exception as e:
        return {"text": "", "error": str(e), "metadata": {}}


# ── XML ───────────────────────────────────────────────────────────────────────
def extract_xml(filepath):
    try:
        from bs4 import BeautifulSoup
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        soup = BeautifulSoup(content, 'lxml-xml')
        text = soup.get_text(separator='\n')
        text = re.sub(r'\n{3,}', '\n\n', text).strip()
        return {
            "text": text,
            "metadata": {},
            "word_count": count_words(text)
        }
    except Exception as e:
        return {"text": "", "error": str(e), "metadata": {}}


# ── MAIN DISPATCHER ───────────────────────────────────────────────────────────
EXTENSION_MAP = {
    '.pdf': extract_pdf,
    '.docx': extract_docx,
    '.doc': extract_docx,
    '.xlsx': extract_xlsx,
    '.xls': extract_xlsx,
    '.pptx': extract_pptx,
    '.ppt': extract_pptx,
    '.csv': extract_csv,
    '.txt': extract_txt,
    '.md': extract_txt,
    '.rtf': extract_txt,
    '.html': extract_html,
    '.htm': extract_html,
    '.json': extract_json,
    '.xml': extract_xml,
    '.log': extract_txt,
    '.py': extract_txt,
    '.js': extract_txt,
    '.ts': extract_txt,
    '.css': extract_txt,
    '.yaml': extract_txt,
    '.yml': extract_txt,
    '.toml': extract_txt,
    '.ini': extract_txt,
    '.conf': extract_txt,
}


def extract_file(filepath):
    """Main entry point - dispatch to correct extractor by extension"""
    ext = Path(filepath).suffix.lower()
    extractor = EXTENSION_MAP.get(ext)
    
    if extractor:
        result = extractor(filepath)
    else:
        # Try as plain text for unknown extensions
        result = extract_txt(filepath)
    
    result['filename'] = os.path.basename(filepath)
    result['extension'] = ext
    result['file_size'] = os.path.getsize(filepath)
    
    if result.get('text'):
        result['summary'] = generate_summary(result['text'])
    
    return result


# ── HELPERS ───────────────────────────────────────────────────────────────────
def count_words(text):
    return len(text.split()) if text else 0


def flatten_json(obj, prefix=''):
    parts = []
    if isinstance(obj, dict):
        for k, v in obj.items():
            parts.extend(flatten_json(v, f"{prefix}{k}: "))
    elif isinstance(obj, list):
        for i, v in enumerate(obj):
            parts.extend(flatten_json(v, prefix))
    else:
        parts.append(f"{prefix}{obj}")
    return "\n".join(parts)


def generate_summary(text, max_sentences=5):
    """Simple extractive summary - no external API needed"""
    if not text or len(text.strip()) < 100:
        return text.strip()
    
    # Clean text
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'--- Page \d+ ---', '', text)
    
    # Split into sentences
    sentences = re.split(r'(?<=[.!?])\s+', text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
    
    if not sentences:
        return text[:500]
    
    if len(sentences) <= max_sentences:
        return ' '.join(sentences)
    
    # Score sentences by position + word frequency
    word_freq = {}
    for sent in sentences:
        for word in sent.lower().split():
            word = re.sub(r'[^a-z]', '', word)
            if len(word) > 3:
                word_freq[word] = word_freq.get(word, 0) + 1
    
    scored = []
    for i, sent in enumerate(sentences):
        score = sum(word_freq.get(re.sub(r'[^a-z]', '', w.lower()), 0)
                    for w in sent.split())
        # Boost first and last sentences
        if i < 3:
            score *= 1.5
        if i >= len(sentences) - 2:
            score *= 1.2
        scored.append((score, i, sent))
    
    top = sorted(scored, reverse=True)[:max_sentences]
    top_sorted = sorted(top, key=lambda x: x[1])
    return ' '.join(s[2] for s in top_sorted)
